# -*- coding: utf-8 -*-
# @Time    : 2020/1/8 10:39
# @Author  : liufeiyafly
import sys
import os
import time
import re
from PyQt5.QtGui import QIcon
from PyQt5.QtWidgets import QApplication, QMainWindow, QFileDialog
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn
from Ui import Ui_MainWindow


class MyMainWindow(QMainWindow, Ui_MainWindow):
    sum_files = 0
    completed_files = 0
    sum_files_actual = 0

    def __init__(self, parent=None):
        super(MyMainWindow, self).__init__(parent)
        self.setupUi(self)
        self.cwd = os.getcwd()  # 获取当前路径
        self.textEdit.acceptDrops()  # 这是textEdit接受拖放
        self.pushButton.setEnabled(False)

        self.setWindowIcon(QIcon(r'./wine.ico'))
        self.pushButton_2.clicked.connect(self.choose_files)
        self.pushButton.clicked.connect(self.start)
        self.textEdit.textChanged.connect(self.pbtn_Enable)
        self.textEdit.textChanged.connect(self.textBrowser2_clear)


    def textBrowser2_clear(self):
        self.textBrowser2.clear()

    def wordsCount(self, string):
        zhongchao_pattern = re.compile(
            '[\u4e00-\u9fa5\uff00-\uffef\u2e80-\u2eff\u3000-\u303f\u31C0-\u31EF]')  # \u2000-\u2046:代表广义标点，包含了中文的引号,但是wordl里并没有统计该字符段
        s_english = re.subn(zhongchao_pattern, ' ',
                            string)  # 要将中文替换为空格，不然有些数字会连起来，再统计就只算一个词了。s_english是去除所有中文剩下的英文。
        Chinese_Words = s_english[1]

        English_Words = re.subn('\s+', ' ', re.subn('[\s\u2000-\u2046]', ' ', s_english[0].strip())[0])[
                            1] + 1  # \s+ :表示至少一个空格

        return (Chinese_Words, English_Words)

    def pbtn_Enable(self):
        self.pushButton.setEnabled(True)  # 只有textedit里面内容变化了，开始按钮才可用。可能不完美，因为主动删除选择的文件后，开始按钮还是可用。

    def choose_files(self):
        # files 是一个list,里面保存了选择文件的文件具体路径。
        files, filetype = QFileDialog.getOpenFileNames(self,
                                                       "选择文件",
                                                       self.cwd,  # 起始路径"C:\\Users\\Administrator\\Desktop"
                                                       "PPTX Files (*.pptx)")  # ;;All Files (*)

        new_files = []  # 里面可能包含用户选择的非pptx文件(一般不会)，要在里面筛选pptx的文件
        for file in files:
            if file[-5:] == '.pptx' or file[-5:] == '.PPTX':
                new_files.append(file)

        if filetype:
            self.textEdit.setText('\n'.join(new_files))

    def get_path_lists(self):
        text = self.textEdit.toPlainText().strip()
        if 'file:///' in text:  # 说明是拖进去的，
            self.sum_files = len(text.split('file:///'))
            files_list = set([f for f in text.split('file:///') if
                              f and os.path.isfile(f) and os.path.getsize(f) > 0])  # set去重。 判断文件存在且大小不为0
        else:  # 说明是选择的文件
            self.sum_files = len(text.split('\n'))
            files_list = set([f for f in text.split('\n') if f and os.path.isfile(f) and os.path.getsize(f) > 0])
        return files_list

    def pptx2docx(self, fileName):  # fileName是“一个”pptx的路径
        auto_object = '<class \'pptx.shapes.autoshape.Shape\'>'
        place_object = '<class \'pptx.shapes.placeholder.SlidePlaceholder\'>'
        group_object = '<class \'pptx.shapes.group.GroupShape\'>'

        def Group_digui(group, l=[]):  # 递归寻找group组合对象里面的文字，（group里面可能套了group）
            for i in group.shapes:
                if str(type(i)) == auto_object or str(type(i)) == place_object:
                    s2 = i.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
                    # print(s2)
                    l.append(s2.strip().replace('\x0b', ' '))
                elif i.shape_type == 6:  # 等于6即 是group组合
                    Group_digui(i)  # 递归访问即可
            return l  # 返回一个列表，里面元素是group内所有的文字。

        prs = Presentation(fileName)
        nums = len(prs.slides)  # 幻灯片的页数

        L = []
        for num in range(nums):
            slide = prs.slides[num]  # 通过索引序号 访问每页幻灯片
            body_shapes = slide.shapes  # slide.shapes 就是每页幻灯片中的所有元素。

            for j in body_shapes:
                if str(type(j)) == auto_object or str(type(j)) == place_object:  # 图形或者文本框是有 .text属性的，可以直接访问
                    s1 = j.text.encode('gbk', 'ignore').decode('gbk', 'ignore')  # 用这个代码解决不能打印和写入文件'gbk'特殊编码的问题了！！！
                    L.append(s1.strip().replace('\x0b', ' '))

                elif str(type(j)) == group_object:
                    text_l = Group_digui(j)
                    L.extend(text_l)

                elif j.shape_type == 19:  # 19代表：表格
                    for row in j.table.rows:
                        for cell in row.cells:
                            s3 = cell.text_frame.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
                            L.append(s3.strip().replace('\x0b', ' '))

        if self.checkbox.isChecked():
            doc = Document()
            doc.styles['Normal'].font.name = u'宋体'
            doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')

            s = ' '.join(L)  # 以空格分开各个内容，也可以用回车
            doc.add_paragraph(s)
            doc.save('{}.docx'.format(fileName[:-5]))  # 生成的docx保存到pptx所在文件夹

            words_count = self.wordsCount(s)
            self.textBrowser2.append('{}:中文字符和朝鲜语单词:{}  非中文单词:{}'.format(os.path.basename(fileName),
                                                                         words_count[0],
                                                                         words_count[1]))
            return words_count[0], words_count[1]
        else:
            s = ' '.join(L)  # 以空格分开各个内容，也可以用回车

            words_count = self.wordsCount(s)
            self.textBrowser2.append('{}:中文字符和朝鲜语单词:{}  非中文单词:{}'.format(os.path.basename(fileName),
                                                                         words_count[0],
                                                                         words_count[1]))
            return words_count[0], words_count[1]

    def main(self, Lists):
        total_Chinese_Words = total_English_Words = 0  # 统计总字数
        for li in Lists:
            Chinese_Words, English_Words = self.pptx2docx(li)
            total_Chinese_Words += Chinese_Words
            total_English_Words += English_Words
            self.completed_files += 1
        return total_Chinese_Words, total_English_Words

    def start(self):
        file_lists = self.get_path_lists()
        if file_lists == None or len(file_lists) == 0:
            return

        self.sum_files_actual = len(file_lists)
        self.textBrowser.setText('完成文件数：{}个\n  总文件数：{}个'.format(self.sum_files_actual, self.sum_files))

        start_time = time.perf_counter()
        total_Chinese_Words, total_English_Words = self.main(file_lists)  # 主要花费时间在这执行步骤
        self.textBrowser2.append('总中文字符和朝鲜语单词:{}，总非中文单词:{}'.format(total_Chinese_Words, total_English_Words))
        total_time = time.perf_counter() - start_time
        self.setWindowTitle('耗时：{:.2f}s'.format(total_time))  # 结束时 标题头 展示耗时。

        self.pushButton.setEnabled(False)  # 运行完要设定为False,以防重复点击‘开始’按钮

if __name__ == '__main__':
    app = QApplication(sys.argv)
    win = MyMainWindow()
    win.show()
    sys.exit(app.exec_())



