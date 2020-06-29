#!/usr/bin/env python3
# -*- coding:utf-8 -*-
'''
@Author: wayne
@Time: 2020/4/17 0017 上午 8:59
'''
import sys
import os
import time
import re
from multi_Ui import Ui_MainWindow
from PyQt5.QtGui import QIcon
from PyQt5.QtCore import Qt
from PyQt5.QtWidgets import QApplication, QMainWindow, QFileDialog, QTableWidgetItem
import threading
from queue import Queue
from pptx import Presentation
from docx import Document
from docx.oxml.ns import qn


# 下面4行用来解决pyqt5生成exe没有图片
# from i_py import img as ico
# import base64
# with open('myi.ico','wb') as f:
#     f.write(base64.b64decode(ico))

class MyMainWindow(QMainWindow, Ui_MainWindow):
    def __init__(self, parent=None):
        super(MyMainWindow, self).__init__(parent)
        self.setupUi(self)
        self.setWindowIcon(QIcon('myi.ico'))
        self.cwd = os.getcwd()
        self.radioButton.setChecked(True)
        self.checkBox.setChecked(True)
        self.pushButton_2.setEnabled(False)
        self.pushButton.clicked.connect(self.choose_files)
        self.pushButton_3.clicked.connect(self.clear_textedit)
        self.textEdit.textChanged.connect(self.pushButton_2Enable)
        self.pushButton_2.clicked.connect(self.start)
        self.textEdit.setPlaceholderText('拖入或选择（多）文件')

        self.checkBox.clicked.connect(self.radio_enable)

        self.queue = Queue()
        self.sum_files = 0
        self.completed_files = 0
        self.AllChineseWords = 0
        self.AllEnglishWords = 0

    def radio_enable(self):
        if not self.checkBox.isChecked():
            self.radioButton.setEnabled(False)
            self.radioButton_2.setEnabled(False)
        else:
            self.radioButton.setEnabled(True)
            self.radioButton_2.setEnabled(True)

    def pushButton_2Enable(self):
        self.pushButton_2.setEnabled(True)

    def clear_textedit(self):
        self.textEdit.clear()
        self.pushButton_2.setEnabled(False)
        self.progressBar.setValue(0)

    def wordsCount(self, string):
        zhongchao_pattern = re.compile(
            '[\u4e00-\u9fa5\uff00-\uffef\u2e80-\u2eff\u3000-\u303f\u31C0-\u31EF]')  # \u2000-\u2046:代表广义标点，包含了中文的引号,但是wordl里并没有统计该字符段
        s_english = re.subn(zhongchao_pattern, ' ',
                            string)  # 要将中文替换为空格，不然有些数字会连起来，再统计就只算一个词了。s_english是去除所有中文剩下的英文。
        Chinese_Words = s_english[1]

        English_Words = re.subn('\s+', ' ', re.subn('[\s\u2000-\u2046]', ' ', s_english[0].strip())[0])[
                            1] + 1  # \s+ :表示至少一个空格

        return Chinese_Words, English_Words

    def choose_files(self):
        # files 是一个list,里面保存了选择文件的文件具体路径。
        files, filetype = QFileDialog.getOpenFileNames(self,
                                                       "选择文件",
                                                       self.cwd,  # 起始路径"C:\\Users\\Administrator\\Desktop"
                                                       "PPTX Files (*.pptx)")  # ;;All Files (*)

        new_files = []  # 里面可能包含用户选择的非pptx文件(一般不会)，要在里面筛选pptx的文件
        for file in files:
            if file[-5:].lower() == '.pptx':
                new_files.append(file)

        if filetype:
            self.textEdit.setText('\n'.join(new_files))
            self.pushButton_2.setEnabled(True)
            self.progressBar.setValue(0)

    def get_path_queue(self):
        text = self.textEdit.toPlainText().strip()
        files_list = []
        if 'file:///' in text:  # 说明是拖进去的，
            self.sum_files = len(text.split('file:///')) - 1
            for f in text.split('file:///'):
                f = f.strip()
                if os.path.isfile(f) and os.path.getsize(f) > 0 and f[-4:].lower() == 'pptx':
                    files_list.append(f)
                else:  # 解决文件名含有 '\xa0' 作为空格的文件。
                    f = os.path.join(os.path.dirname(f), os.path.basename(f).replace(' ', '\xa0'))
                    if os.path.isfile(f) and os.path.getsize(f) > 0 and f[-4:].lower() == 'pptx':
                        files_list.append(f)
        else:  # 否则说明是选择的文件
            self.sum_files = len(text.split('\n'))
            for f in text.split('\n'):
                if os.path.isfile(f) and os.path.getsize(f) > 0:
                    files_list.append(f)
                else:  # 解决文件名含有 '\xa0' 作为空格的文件。
                    f = os.path.join(os.path.dirname(f), os.path.basename(f).replace(' ', '\xa0'))
                    if os.path.isfile(f) and os.path.getsize(f) > 0:
                        files_list.append(f)
        for i in set(files_list):
            self.queue.put_nowait(i)
        return self.queue

    def pptx2docx(self, file_q):  # fileName是“一个”pptx的路径

        auto_object = '<class \'pptx.shapes.autoshape.Shape\'>'
        place_object = '<class \'pptx.shapes.placeholder.SlidePlaceholder\'>'
        group_object = '<class \'pptx.shapes.group.GroupShape\'>'

        def Group_digui(group, l=[]):  # 递归寻找group组合对象里面的文字，（group里面可能套了group）

            for i in group.shapes:
                if str(type(i)) == auto_object or str(type(i)) == place_object:
                    s2 = i.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
                    # print(s2)
                    l.append(s2.strip().replace('\x0b', ' '))
                elif i.shape_type == 6:  # 等于6即 是group组合
                    l.extend(Group_digui(i, []))  # 递归访问即可

            return l  # 返回一个列表，里面元素是group内所有的文字。

        while not file_q.empty():
            fileName = file_q.get_nowait()
            try:
                prs = Presentation(fileName)
            except Exception as e:
                pass
            else:
                nums = len(prs.slides)  # 幻灯片的页数

                L = []
                for num in range(nums):
                    slide = prs.slides[num]  # 通过索引序号 访问每页幻灯片
                    body_shapes = slide.shapes  # slide.shapes 就是每页幻灯片中的所有元素。

                    for j in body_shapes:
                        if str(type(j)) == auto_object or str(type(j)) == place_object:  # 图形或者文本框是有 .text属性的，可以直接访问
                            s1 = j.text.encode('gbk', 'ignore').decode('gbk',
                                                                       'ignore')  # 用这个代码解决不能打印和写入文件'gbk'特殊编码的问题了！！！
                            L.append(s1.strip().replace('\x0b', ' '))

                        elif str(type(j)) == group_object:
                            text_l = Group_digui(j, [])  # Group_digui()的第二个参数必须加，且为空列表[]
                            L.extend(text_l)

                        elif j.shape_type == 19:  # 19代表：表格
                            for row in j.table.rows:
                                for cell in row.cells:
                                    s3 = cell.text_frame.text.encode('gbk', 'ignore').decode('gbk', 'ignore')
                                    L.append(s3.strip().replace('\x0b', ' '))
                all_words_with_space = ' '.join(L)  # 以空格分开各个内容

                if self.checkBox.isChecked():
                    doc = Document()
                    doc.styles['Normal'].font.name = u'宋体'
                    doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')
                    if self.radioButton.isChecked():
                        s = all_words_with_space
                    elif self.radioButton_2.isChecked():
                        s = '\n'.join(L)  # 以回车分开各个内容
                    doc.add_paragraph(s)
                    doc.save('{}.docx'.format(fileName[:-5]))  # 生成的docx保存到pptx所在文件夹

                words_count = self.wordsCount(all_words_with_space)  # 用空格分隔的总内容，来统计字数。
                # self.lock.acquire()
                item = QTableWidgetItem(os.path.basename(fileName))
                item.setTextAlignment(Qt.AlignLeft)
                self.tableWidget_2.setItem(self.completed_files, 0, item)

                item = QTableWidgetItem(str(words_count[0]))
                item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
                self.tableWidget_2.setItem(self.completed_files, 1, item)

                item = QTableWidgetItem(str(words_count[1]))
                item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
                self.tableWidget_2.setItem(self.completed_files, 2, item)

                self.AllChineseWords += words_count[0]
                self.AllEnglishWords += words_count[1]

                self.completed_files += 1
                self.progressBar.setValue(self.completed_files)

    def start(self):
        file_queue = self.get_path_queue()
        if file_queue.empty():
            return
        print(file_queue.queue)
        self.sum_files = file_queue.qsize()
        self.completed_files = 0
        self.AllChineseWords = 0
        self.AllEnglishWords = 0

        self.tableWidget_2.setRowCount(self.sum_files + 1)

        self.progressBar.setMaximum(self.sum_files)
        start_time = time.perf_counter()
        if self.sum_files > 1:
            t1 = threading.Thread(target=self.pptx2docx, args=(file_queue,))
            t2 = threading.Thread(target=self.pptx2docx, args=(file_queue,))
            t1.start()
            t2.start()
            t1.join()  # 等子线程结束之后 主线程才结束。才可统计时间。
            t2.join()
        else:
            self.pptx2docx(file_queue)

        total_time = time.perf_counter() - start_time

        item = QTableWidgetItem(str(self.sum_files))
        item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
        self.tableWidget.setItem(0, 0, item)

        item = QTableWidgetItem(str(self.completed_files))
        item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
        self.tableWidget.setItem(0, 1, item)

        item = QTableWidgetItem('总计：')
        item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
        self.tableWidget_2.setItem(self.completed_files, 0, item)

        item = QTableWidgetItem(str(self.AllChineseWords))
        item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
        self.tableWidget_2.setItem(self.completed_files, 1, item)

        item = QTableWidgetItem(str(self.AllEnglishWords))
        item.setTextAlignment(Qt.AlignRight | Qt.AlignCenter)
        self.tableWidget_2.setItem(self.completed_files, 2, item)

        self.setWindowTitle('耗时：{:.2f}s'.format(total_time))  # 结束时 标题头 展示耗时。
        self.pushButton_2.setEnabled(False)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    win = MyMainWindow()
    win.show()
    sys.exit(app.exec_())
